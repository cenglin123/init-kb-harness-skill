#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
extract_office.py — office 文档纯文本提取（sidecar 机制）

把 vault 内的 .docx / .xlsx / .pptx / .pdf（OFFICE_EXTRACT_EXTS 可配）提取为
.meta/office-extracts/<镜像路径>.md sidecar，随后由 embed / bm25 / build_index
纳入检索链路。检索命中 sidecar 时，ask.py 会显示指向源 office 文件。

- 增量：源文件 content_hash 未变则跳过（--full 强制重提）
- 并发：ThreadPoolExecutor（MAINTAIN_CONCURRENCY，默认 6）
- 孤儿清理：源文件删除/改名后对应 sidecar 自动删除
- 老格式 .doc/.xls/.ppt 不提取，清单写入 _legacy-formats.md 提示用户转存新格式
- 隐私目录（category-privacy.md / PRIVACY_DIRS）内文件不扫描、不提取

用法:
    python .meta/scripts/extract_office.py           # 增量
    python .meta/scripts/extract_office.py --full    # 全量重提
"""

import sys
import argparse
from concurrent.futures import ThreadPoolExecutor, as_completed
from datetime import datetime
from pathlib import Path

sys.path.insert(0, str(Path(__file__).parent))
from common import (
    VAULT_ROOT, rel_path, ensure_parent, content_hash,
    OFFICE_EXTRACT_EXTS, OFFICE_EXTRACTS_DIR, MAINTAIN_CONCURRENCY, ENV,
    scan_office_docs, scan_office_legacy_docs, scan_office_extracts,
)

# 单文件提取文本上限（防超大表格灌爆索引；0 = 不限）
MAX_CHARS = int(ENV.get('OFFICE_EXTRACT_MAX_CHARS', '200000'))
# xlsx 每 sheet 最大行数（0 = 不限）
MAX_ROWS_PER_SHEET = int(ENV.get('OFFICE_EXTRACT_MAX_ROWS', '2000'))

LEGACY_REPORT = OFFICE_EXTRACTS_DIR / '_legacy-formats.md'


# ─── 各格式提取器（缺依赖时该格式整体跳过并提示）──────────────────────────

def _extract_docx(path: Path) -> str:
    import docx  # python-docx
    doc = docx.Document(str(path))
    parts = []
    for para in doc.paragraphs:
        t = para.text.strip()
        if t:
            parts.append(t)
    for table in doc.tables:
        for row in table.rows:
            cells = [c.text.strip().replace('\n', ' ') for c in row.cells]
            if any(cells):
                parts.append(' | '.join(cells))
    return '\n'.join(parts)


def _extract_xlsx(path: Path) -> str:
    import openpyxl
    wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
    parts = []
    try:
        for ws in wb.worksheets:
            parts.append(f'## Sheet: {ws.title}')
            for i, row in enumerate(ws.iter_rows(values_only=True)):
                if MAX_ROWS_PER_SHEET and i >= MAX_ROWS_PER_SHEET:
                    parts.append(f'（… 已截断，仅保留前 {MAX_ROWS_PER_SHEET} 行）')
                    break
                cells = ['' if v is None else str(v).strip() for v in row]
                if any(cells):
                    parts.append(' | '.join(cells))
            parts.append('')
    finally:
        wb.close()
    return '\n'.join(parts)


def _extract_pptx(path: Path) -> str:
    import pptx  # python-pptx
    prs = pptx.Presentation(str(path))
    parts = []
    for i, slide in enumerate(prs.slides, 1):
        parts.append(f'## Slide {i}')
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = ''.join(run.text for run in para.runs).strip()
                    if t:
                        parts.append(t)
            if getattr(shape, 'has_table', False) and shape.has_table:
                for row in shape.table.rows:
                    cells = [c.text.strip().replace('\n', ' ') for c in row.cells]
                    if any(cells):
                        parts.append(' | '.join(cells))
        try:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                parts.append(f'> 备注: {notes}')
        except Exception:
            pass
        parts.append('')
    return '\n'.join(parts)


def _extract_pdf(path: Path) -> str:
    from pypdf import PdfReader
    reader = PdfReader(str(path))
    parts = []
    for i, page in enumerate(reader.pages, 1):
        try:
            text = (page.extract_text() or '').strip()
        except Exception:
            text = ''
        if text:
            parts.append(f'## Page {i}')
            parts.append(text)
            parts.append('')
    return '\n'.join(parts)


EXTRACTORS = {
    '.docx': ('python-docx', _extract_docx),
    '.xlsx': ('openpyxl', _extract_xlsx),
    '.pptx': ('python-pptx', _extract_pptx),
    '.pdf': ('pypdf', _extract_pdf),
}


# ─── sidecar 读写 ──────────────────────────────────────────────────────────

def sidecar_path(src: Path) -> Path:
    """源文件 → .meta/office-extracts/<镜像路径>.md"""
    rel = src.relative_to(VAULT_ROOT)
    return OFFICE_EXTRACTS_DIR / (str(rel) + '.md')


def sidecar_hash(sc: Path) -> str:
    """读 sidecar frontmatter 里记录的 source_hash；失败返回空串。"""
    try:
        for line in sc.read_text(encoding='utf-8').splitlines()[:10]:
            if line.startswith('source_hash:'):
                return line.split(':', 1)[1].strip()
    except Exception:
        pass
    return ''


def write_sidecar(src_rel: str, h: str, extractor: str, text: str, sc: Path):
    truncated = False
    if MAX_CHARS and len(text) > MAX_CHARS:
        text = text[:MAX_CHARS]
        truncated = True
    fm = (
        '---\n'
        f'source: {src_rel}\n'
        f'source_hash: {h}\n'
        f'extracted_at: {datetime.now().isoformat(timespec="seconds")}\n'
        f'extractor: {extractor}\n'
        'type: office-extract\n'
        '---\n\n'
        f'# {Path(src_rel).name}\n\n'
        '<!-- 由 extract_office.py 自动提取，请勿手动编辑；编辑请改源文件 -->\n\n'
    )
    body = text if text.strip() else '（未提取到文本内容）'
    if truncated:
        body += f'\n\n（… 已截断，仅保留前 {MAX_CHARS} 字符）'
    ensure_parent(sc)
    sc.write_text(fm + body + '\n', encoding='utf-8', newline='')


def _extract_one(src: Path, full: bool):
    """单文件提取。返回 (rel, status, message)；status ∈ {ok, skip, fail}。"""
    rel = rel_path(src)
    ext = src.suffix.lower()
    lib, fn = EXTRACTORS[ext]
    sc = sidecar_path(src)
    h = content_hash(src)
    if not full and sc.exists() and sidecar_hash(sc) == h:
        return rel, 'skip', ''
    try:
        text = fn(src)
    except ImportError:
        return rel, 'fail', f'缺依赖 {lib}（pip install {lib}）'
    except Exception as e:
        return rel, 'fail', f'{type(e).__name__}: {e}'
    write_sidecar(rel, h, lib, text, sc)
    return rel, 'ok', f'{len(text)} 字符'


# ─── 孤儿清理 / 老格式登记 ────────────────────────────────────────────────

def cleanup_orphans(valid_sidecars: set) -> int:
    """删除源文件已不存在的 sidecar。"""
    removed = 0
    for sc in list(scan_office_extracts()):
        if sc not in valid_sidecars:
            try:
                sc.unlink()
                removed += 1
            except Exception:
                pass
    return removed


def report_legacy():
    legacy = sorted(rel_path(f) for f in scan_office_legacy_docs())
    if not legacy:
        if LEGACY_REPORT.exists():
            LEGACY_REPORT.unlink()
        return legacy
    lines = [
        '# 老格式 office 文件（无法提取）',
        '',
        '> 以下 .doc/.xls/.ppt 为老二进制格式，本管线不提取。',
        '> 请用 Office/WPS 转存为 .docx/.xlsx/.pptx 后自动纳入检索。',
        '',
    ]
    lines += [f'- `{p}`' for p in legacy]
    ensure_parent(LEGACY_REPORT)
    LEGACY_REPORT.write_text('\n'.join(lines) + '\n', encoding='utf-8', newline='')
    return legacy


# ─── 主流程 ────────────────────────────────────────────────────────────────

def main(full=False):
    all_docs = sorted(scan_office_docs(), key=rel_path)
    docs = list(all_docs)
    print(f"  office 文档: {len(docs)} 个（扩展名 {', '.join(sorted(OFFICE_EXTRACT_EXTS))}）")

    # 缺依赖的格式整体提示（提前探测，避免每文件重复报错）
    missing = []
    present_exts = {d.suffix.lower() for d in docs}
    for ext in sorted(present_exts):
        lib, _ = EXTRACTORS[ext]
        try:
            __import__({'python-docx': 'docx', 'openpyxl': 'openpyxl',
                        'python-pptx': 'pptx', 'pypdf': 'pypdf'}[lib])
        except ImportError:
            missing.append((ext, lib))
    if missing:
        for ext, lib in missing:
            print(f"  ⚠️  {ext} 提取依赖缺失：pip install {lib}（该格式本轮跳过）")
        skip_exts = {ext for ext, _ in missing}
        docs = [d for d in docs if d.suffix.lower() not in skip_exts]

    ok = skipped = failed = 0
    if docs:
        with ThreadPoolExecutor(max_workers=MAINTAIN_CONCURRENCY) as pool:
            futures = {pool.submit(_extract_one, d, full): d for d in docs}
            for fut in as_completed(futures):
                rel, status, msg = fut.result()
                if status == 'ok':
                    ok += 1
                    print(f"  ✓ {rel} ({msg})")
                elif status == 'skip':
                    skipped += 1
                else:
                    failed += 1
                    print(f"  ✗ {rel}: {msg}")

    # 孤儿判定基于全部扫描到的源文件（含因缺依赖本轮跳过的格式）——
    # 只有源文件真正消失才算孤儿，依赖缺失不触发误删
    valid = {sidecar_path(d) for d in all_docs}
    removed = cleanup_orphans(valid)
    legacy = report_legacy()

    print(f"  ✓ office 提取完成：{ok} 提取 / {skipped} 未变跳过 / {failed} 失败 / {removed} 孤儿清理")
    if legacy:
        print(f"  ⚠️  {len(legacy)} 个老格式文件无法提取（见 .meta/office-extracts/_legacy-formats.md）")
    return 0


if __name__ == '__main__':
    # --- host guard ---
    from common import is_primary_host
    if not is_primary_host():
        print("FATAL: this script must run on PRIMARY_HOST", file=sys.stderr)
        sys.exit(1)
    # --- /host guard ---
    parser = argparse.ArgumentParser(description='office 文档纯文本提取（sidecar）')
    parser.add_argument('--full', action='store_true', help='忽略 hash 全量重提')
    args = parser.parse_args()
    sys.exit(main(full=args.full))
